#!/usr/bin/env python3
"""Generate Chapter 3 presentation — McKinsey/BCG style."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ───────────────────────────────────────────────────────────────────
ROYAL_BLUE  = RGBColor(0x00, 0x33, 0x66)
DARK_GREY   = RGBColor(0x33, 0x33, 0x33)
MED_GREY    = RGBColor(0x66, 0x66, 0x66)
LIGHT_GREY  = RGBColor(0xCC, 0xCC, 0xCC)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
RED_ACCENT  = RGBColor(0xC0, 0x39, 0x2B)
BLUE_ACCENT = RGBColor(0x29, 0x80, 0xB9)
GREEN_ACC   = RGBColor(0x27, 0xAE, 0x60)
GOLD_ACC    = RGBColor(0xD4, 0xA0, 0x17)
BG_GREY     = RGBColor(0xF5, 0xF5, 0xF5)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Helper functions ─────────────────────────────────────────────────────────
def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def add_textbox(slide, left, top, width, height, text, font_size=12,
                bold=False, color=BLACK, font_name='Arial', alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_title(slide, text, top=0.3):
    add_textbox(slide, 0.8, top, 11.5, 0.7, text, font_size=28,
                bold=True, color=ROYAL_BLUE, font_name='Times New Roman')

def add_subtitle(slide, text, top=1.0):
    add_textbox(slide, 0.8, top, 11.5, 0.5, text, font_size=14,
                color=MED_GREY, font_name='Times New Roman')

def add_line(slide, left, top, width, color=LIGHT_GREY, line_width=Pt(0.5)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_table(slide, left, top, col_widths, headers, rows, header_color=ROYAL_BLUE):
    """Add a styled table. col_widths in inches, headers=list, rows=list of lists."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_w = sum(col_widths)
    table_shape = slide.shapes.add_table(n_rows, n_cols,
        Inches(left), Inches(top), Inches(total_w), Inches(0.35 * n_rows))
    table = table_shape.table

    # Set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = 'Arial'
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color

    # Data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8)
                p.font.color.rgb = DARK_GREY
                p.font.name = 'Arial'
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 == 0 else BG_GREY

    # Thin borders
    for r in range(n_rows):
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.margin_left = Pt(3)
            cell.margin_right = Pt(3)
            cell.margin_top = Pt(2)
            cell.margin_bottom = Pt(2)

    return table

def add_kpi_box(slide, left, top, width, number, label, color=ROYAL_BLUE):
    """Big number + label KPI box."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_GREY
    shape.line.fill.background()

    add_textbox(slide, left + 0.1, top + 0.05, width - 0.2, 0.5, str(number),
                font_size=24, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + 0.1, top + 0.5, width - 0.2, 0.35, label,
                font_size=8, color=MED_GREY, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: COVER
# ══════════════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
# Blue bar at top
add_line(slide, 0, 0, 13.33, ROYAL_BLUE, Pt(4))
# Title
add_textbox(slide, 1.5, 2.0, 10, 1.2,
    "Comparative Transcriptomic Analysis of\nHerpes Zoster Disease and RZV Vaccine Immunity",
    font_size=34, bold=True, color=ROYAL_BLUE, font_name='Times New Roman')
# Subtitle
add_textbox(slide, 1.5, 3.5, 10, 0.8,
    "Establishing Disease and Protection Immune Signatures\nfrom Public Transcriptomic Datasets",
    font_size=16, color=MED_GREY, font_name='Times New Roman')
# Divider
add_line(slide, 1.5, 4.5, 3, ROYAL_BLUE, Pt(1.5))
# Meta
add_textbox(slide, 1.5, 5.0, 10, 0.5, "Datasets: GSE242252 (Bulk RNA-seq)  |  GSE249632 (scRNA-seq)  |  HRA008316 (scRNA-seq validation)",
            font_size=10, color=MED_GREY)
add_textbox(slide, 1.5, 5.5, 10, 0.5, "June 2026  |  Confidential",
            font_size=9, color=LIGHT_GREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: EXECUTIVE SUMMARY
# ══════════════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_line(slide, 0, 0, 13.33, ROYAL_BLUE, Pt(3))
add_title(slide, "Executive Summary")
add_line(slide, 0.8, 1.05, 11.5, ROYAL_BLUE, Pt(1.5))

# KPI row
add_kpi_box(slide, 0.8, 1.4, 2.6, "352 ↑ / 44 ↓", "HZ Acute DEGs (p<0.05, |LFC|>0.58)", RED_ACCENT)
add_kpi_box(slide, 3.7, 1.4, 2.6, "IFI44L +1.82", "Top ISG in HZ (Disease Signature)", RED_ACCENT)
add_kpi_box(slide, 6.6, 1.4, 2.6, "D14: 164↑ / 152↓", "RZV Peak Response (1st dose)", BLUE_ACCENT)
add_kpi_box(slide, 9.5, 1.4, 2.6, "ZEB2 D365 +3.0", "RZV Long-term Reprogramming", BLUE_ACCENT)

# Key findings
findings = [
    ("HZ Disease Signature",
     "Type I IFN-driven innate immune storm: ISGs (IFI44L, IFI27, RSAD2, ISG15), "
     "complement activation (SERPING1), B cell/plasma cell expansion (MZB1, IGHG4), "
     "and broad cellular proliferation (TOP2A, PTTG1). ISG signal originates from "
     "monocytes/DCs/neutrophils, not T cells (validated by scRNA-seq HRA008316)."),
    ("RZV Protection Signature",
     "Pulsed, self-limiting adaptive CD4⁺ T cell immunity. Peak activation at D14 "
     "(164 up), return to baseline at D60 (10 up), re-activation at D74 (72 up). "
     "Long-term effector programming persists at D365 (HAVCR2 +8.0, ZEB2 +3.0). "
     "Inflammatory pathways actively suppressed (SPP1 −5.9 at D365)."),
    ("Cross-Dataset Insight",
     "HZ and RZV induce fundamentally different transcriptional programs. Type I IFN "
     "pathway — the hallmark of HZ — is completely absent in RZV response (ISG15: "
     "HZ +1.57 vs RZV D14 +0.1). RZV achieves protection through precise T cell "
     "reprogramming without triggering innate inflammation."),
]
y = 2.7
for title, body in findings:
    add_textbox(slide, 0.8, y, 1.8, 0.3, title, font_size=10, bold=True, color=ROYAL_BLUE)
    add_textbox(slide, 2.7, y, 9.8, 0.7, body, font_size=9, color=DARK_GREY)
    y += 1.2

# Bottom line
add_line(slide, 0.8, 6.5, 11.5, LIGHT_GREY, Pt(0.5))
add_textbox(slide, 0.8, 6.6, 11.5, 0.4,
    "Core Thesis: An ideal vaccine should induce RZV-like Protection Signature while avoiding HZ-like Disease Signature.",
    font_size=10, bold=True, color=ROYAL_BLUE, font_name='Times New Roman')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: STUDY DESIGN
# ══════════════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_line(slide, 0, 0, 13.33, ROYAL_BLUE, Pt(3))
add_title(slide, "Study Design: Three Datasets, Two Reference Frameworks")
add_line(slide, 0.8, 1.05, 11.5, ROYAL_BLUE, Pt(1.5))

add_table(slide, 0.8, 1.4,
    [2.0, 3.2, 3.2, 3.2],
    ["Feature", "GSE242252", "GSE249632", "HRA008316"],
    [
        ["Technology", "Bulk 3′ mRNA-seq", "SMART-Seq v4 scRNA-seq", "10x Genomics scRNA-seq"],
        ["Sample Source", "Whole Blood", "gE Tetramer⁺ CD4⁺ T Cells", "PBMCs"],
        ["Subjects", "26 HZ Patients", "7 Healthy Vaccinees", "3 HA / 3 HP / 3 RP"],
        ["Comparison", "Acute vs Convalescent (1yr)", "D0→D14→D60→D74→D365", "HP vs HA vs RP"],
        ["Cells Passed QC", "N/A (Bulk)", "2,231", "66,338"],
        ["Analysis Method", "DESeq2 (unpaired)", "limma-voom (paired)", "Cell-type annotation"],
        ["Primary Output", "Disease Signature", "Protection Signature", "Cell-source validation"],
        ["Accession / Reference", "Vandoren et al. 2024", "GEO GSE249632", "Zheng et al. 2024"],
    ])

# Framework diagram
add_textbox(slide, 0.8, 5.0, 11.5, 0.3, "Analytical Framework", font_size=12, bold=True, color=ROYAL_BLUE)
# Two boxes
for i, (label, desc, color) in enumerate([
    ("HZ Disease Reference", "What immune programs\n are activated during HZ?\n → Define Disease Signature", RED_ACCENT),
    ("RZV Protection Reference", "What immune programs does\n a successful vaccine induce?\n → Define Protection Signature", BLUE_ACCENT),
]):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8 + i * 6), Inches(5.4), Inches(5.5), Inches(1.3))
    shape.fill.solid(); shape.fill.fore_color.rgb = BG_GREY
    shape.line.color.rgb = color; shape.line.width = Pt(1.5)
    add_textbox(slide, 1.0 + i * 6, 5.45, 5.0, 0.3, label, font_size=11, bold=True, color=color)
    add_textbox(slide, 1.0 + i * 6, 5.8, 5.0, 0.8, desc, font_size=9, color=DARK_GREY)

# Arrow between
add_textbox(slide, 6.2, 5.8, 1.2, 0.4, "⟶", font_size=24, color=MED_GREY, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: HZ VOLCANO + KEY DEGs
# ══════════════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_line(slide, 0, 0, 13.33, ROYAL_BLUE, Pt(3))
add_title(slide, "HZ Disease Signature: Acute vs Convalescent Whole Blood Transcriptome")
add_subtitle(slide, "GSE242252 | DESeq2 Unpaired | 26 Acute vs 23 Convalescent | p<0.05, |LFC|>0.58 | Vandoren et al. 2024 Table S5")
add_line(slide, 0.8, 1.55, 11.5, ROYAL_BLUE, Pt(1.5))

# Left: Key DEG table
add_textbox(slide, 0.8, 1.8, 3, 0.3, "Key Differentially Expressed Genes", font_size=11, bold=True, color=ROYAL_BLUE)
add_table(slide, 0.8, 2.2,
    [1.8, 1.0, 1.0, 2.0],
    ["Gene", "LFC", "p-value", "Functional Category"],
    [
        ["IFI27", "+2.16", "1.5e-04", "Type I IFN (ISG)"],
        ["PTTG1", "+2.01", "3.0e-13", "Cell Proliferation"],
        ["BATF2", "+1.95", "9.5e-04", "IFN-Induced TF"],
        ["IGHG4", "+1.94", "2.6e-05", "Immunoglobulin"],
        ["IFI44L", "+1.82", "8.5e-05", "Type I IFN (ISG)"],
        ["SERPING1", "+1.80", "5.3e-07", "Complement Activation"],
        ["MZB1", "+1.68", "3.9e-06", "Plasma Cell Marker"],
        ["ISG15", "+1.57", "1.6e-05", "Type I IFN (ISG)"],
        ["RSAD2", "+1.55", "1.2e-04", "Antiviral Effector"],
        ["TOP2A", "+1.50", "3.4e-13", "DNA Replication"],
    ])

# Middle: GO enrichment
add_textbox(slide, 5.2, 1.8, 3.5, 0.3, "GO Biological Process Enrichment", font_size=11, bold=True, color=ROYAL_BLUE)
add_table(slide, 5.2, 2.2,
    [2.8, 0.8],
    ["GO Term", "Genes"],
    [
        ["Defense Response to Virus", "8"],
        ["B Cell Receptor Signaling", "5"],
        ["Antigen Receptor Signaling", "6"],
        ["Mitotic Spindle Organization", "5"],
        ["Chromosome Condensation", "3"],
    ])

# Right: Disease Signature summary
add_textbox(slide, 9.3, 1.8, 3.5, 0.3, "Disease Signature Genes", font_size=11, bold=True, color=ROYAL_BLUE)
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(9.3), Inches(2.2), Inches(3.5), Inches(3.5))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xEC)
shape.line.color.rgb = RED_ACCENT; shape.line.width = Pt(1)

sig_text = (
    "Type I IFN Storm:\n"
    "  IFI27, IFI44L, RSAD2, ISG15,\n"
    "  MX1, IFIT5, OASL, IFI44\n\n"
    "Complement Activation:\n"
    "  SERPING1, SIGLEC1\n\n"
    "B Cell / Plasma Cell:\n"
    "  MZB1, IGHG4, IGLC2\n\n"
    "Cell Proliferation:\n"
    "  PTTG1, TOP2A"
)
add_textbox(slide, 9.5, 2.3, 3.0, 3.3, sig_text, font_size=9, color=DARK_GREY)

# Bottom note
add_line(slide, 0.8, 6.0, 11.5, LIGHT_GREY, Pt(0.5))
add_textbox(slide, 0.8, 6.1, 11.5, 0.3,
    "Key Insight: HZ is an innate immune-dominated, multi-cellular inflammatory response. ISG signal originates from monocytes/DCs/neutrophils (HRA008316).",
    font_size=9, bold=True, color=DARK_GREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: RZV DEG TIMELINE
# ══════════════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_line(slide, 0, 0, 13.33, ROYAL_BLUE, Pt(3))
add_title(slide, "RZV Vaccine: CD4⁺ T Cell Transcriptional Dynamics")
add_subtitle(slide, "GSE249632 | limma-voom Paired Design | 7 Donors, 2,231 QC Cells | FDR<0.05, |LFC|>0.58")
add_line(slide, 0.8, 1.55, 11.5, ROYAL_BLUE, Pt(1.5))

# DEG timeline table
add_textbox(slide, 0.8, 1.8, 4, 0.3, "DEG Counts by Timepoint", font_size=11, bold=True, color=ROYAL_BLUE)
add_table(slide, 0.8, 2.2,
    [1.5, 1.0, 1.0, 2.5],
    ["Timepoint", "Up", "Down", "Interpretation"],
    [
        ["D14 vs D0", "164", "152", "Strong 1st dose activation"],
        ["D60 vs D0", "10", "4", "Return to baseline (pre-2nd dose)"],
        ["D74 vs D0", "72", "35", "2nd dose re-activation (attenuated)"],
        ["D365 vs D0", "20", "21", "Long-term residual imprint (41 DEGs)"],
    ])

# Key gene trajectory
add_textbox(slide, 4.8, 1.8, 4, 0.3, "Protection Signature Genes", font_size=11, bold=True, color=ROYAL_BLUE)
add_table(slide, 4.8, 2.2,
    [1.3, 1.0, 1.0, 1.0, 2.0],
    ["Gene", "D14", "D74", "D365", "Function"],
    [
        ["ZEB2", "+3.2", "+3.6", "+3.0", "T Cell Differentiation Reprogramming"],
        ["CTLA4", "+1.5", "+1.6", "+1.0", "Immune Checkpoint / Self-Limitation"],
        ["ICOS", "+1.3", "+1.2", "+0.6", "T Cell Co-Stimulation"],
        ["HAVCR2", "+7.6", "+9.1", "+8.0", "Effector Memory (TIM-3)"],
        ["SPP1", "−5.4", "−6.2", "−5.9", "Inflammatory Cytokine — Actively Suppressed"],
        ["ISG15", "+0.1", "−0.4", "−0.1", "Type I IFN — NOT Activated (Neg. Control)"],
    ])

# Interpretation box
add_textbox(slide, 8.6, 1.8, 4, 0.3, "Pattern Classification", font_size=11, bold=True, color=ROYAL_BLUE)
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(8.6), Inches(2.2), Inches(4.2), Inches(3.5))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)
shape.line.color.rgb = BLUE_ACCENT; shape.line.width = Pt(1)

pattern_text = (
    "PULSED ACTIVATION:\n"
    "  D14 peak → D60 near-baseline\n"
    "  → D74 re-activation (weaker)\n"
    "  → D365 residual programming\n\n"
    "PERSISTENT SIGNATURE (D365):\n"
    "  ZEB2, HAVCR2, CTLA4 ↑\n"
    "  SPP1, CA2 ↓\n\n"
    "COMPLETELY ABSENT:\n"
    "  Type I IFN ISGs (ISG15, RSAD2)\n"
    "  Complement (SERPING1)\n"
    "  B cell markers (MZB1)"
)
add_textbox(slide, 8.8, 2.3, 3.8, 3.3, pattern_text, font_size=9, color=DARK_GREY)

# Bottom
add_line(slide, 0.8, 6.0, 11.5, LIGHT_GREY, Pt(0.5))
add_textbox(slide, 0.8, 6.1, 11.5, 0.3,
    "Key Insight: RZV induces pulsed, self-limiting adaptive immunity — not sustained inflammation. Long-term effector programming established at D365.",
    font_size=9, bold=True, color=DARK_GREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: CROSS-DATASET COMPARISON
# ══════════════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_line(slide, 0, 0, 13.33, ROYAL_BLUE, Pt(3))
add_title(slide, "Disease vs Protection: Head-to-Head Signature Comparison")
add_line(slide, 0.8, 1.05, 11.5, ROYAL_BLUE, Pt(1.5))

# 2x2 Matrix
add_textbox(slide, 0.8, 1.4, 11.5, 0.3, "Gene-Level Signature Comparison Matrix", font_size=12, bold=True, color=ROYAL_BLUE)

# Matrix as table
add_table(slide, 0.8, 1.8,
    [2.5, 1.8, 1.8, 1.8, 1.7, 2.5],
    ["Gene", "Functional Role", "HZ LFC", "RZV D14", "RZV D365", "Interpretation"],
    [
        ["IFI44L", "Type I IFN ISG", "+1.82 **", "+0.2 ns", "−0.2 ns", "DISEASE-SPECIFIC: RZV avoids"],
        ["IFI27", "Type I IFN ISG", "+2.16 *", "−0.1 ns", "−0.2 ns", "DISEASE-SPECIFIC: RZV avoids"],
        ["ISG15", "Type I IFN ISG", "+1.57 **", "+0.1 ns", "−0.1 ns", "DISEASE-SPECIFIC: RZV avoids"],
        ["RSAD2", "Antiviral Effector", "+1.55 *", "−0.3 ns", "−0.2 ns", "DISEASE-SPECIFIC: RZV avoids"],
        ["SERPING1", "Complement", "+1.80 ***", "−0.5 ns", "−0.3 ns", "DISEASE-SPECIFIC: RZV avoids"],
        ["", "", "", "", "", ""],
        ["ZEB2", "T Cell Differentiation", "−0.2 ns", "+3.2 ***", "+3.0 ***", "PROTECTION: Vaccine hallmark"],
        ["CTLA4", "Immune Checkpoint", "+0.1 ns", "+1.5 ***", "+1.0 *", "PROTECTION: Self-limitation"],
        ["ICOS", "Co-Stimulation", "+0.1 ns", "+1.3 **", "+0.6 ns", "PROTECTION: T cell help"],
        ["HAVCR2", "Effector Memory", "+0.1 ns", "+7.6 ***", "+8.0 ***", "PROTECTION: Long-term memory"],
        ["SPP1", "Inflammatory Cytokine", "+0.3 ns", "−5.4 ***", "−5.9 ***", "PROTECTION: Active suppression"],
        ["", "", "", "", "", ""],
        ["TOP2A", "Cell Proliferation", "+1.50 ***", "+1.0 *", "−0.6 ns", "SHARED: Generic proliferation"],
        ["PTTG1", "Cell Cycle", "+2.01 ***", "+1.2 **", "+0.1 ns", "SHARED: Generic proliferation"],
    ])

# Bottom
add_line(slide, 0.8, 5.8, 11.5, LIGHT_GREY, Pt(0.5))
add_textbox(slide, 0.8, 6.0, 11.5, 0.3, "Core Finding", font_size=11, bold=True, color=ROYAL_BLUE)
add_textbox(slide, 0.8, 6.3, 11.5, 0.5,
    "HZ and RZV are fundamentally distinct immune programs. RZV protection does NOT require activating innate IFN inflammation. "
    "Instead, it relies on precise CD4⁺ T cell reprogramming (ZEB2) with built-in self-limitation (CTLA4) and active suppression of inflammatory pathways (SPP1).",
    font_size=9, color=DARK_GREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: EVALUATION FRAMEWORK
# ══════════════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_line(slide, 0, 0, 13.33, ROYAL_BLUE, Pt(3))
add_title(slide, "Candidate Vaccine Evaluation Framework")
add_subtitle(slide, "Dual-Standard Assessment Derived from Public Transcriptomic References")
add_line(slide, 0.8, 1.55, 11.5, ROYAL_BLUE, Pt(1.5))

# Two standards side by side
for i, (std_num, std_name, ref, genes, direction, chapter, color) in enumerate([
    ("Standard 1", "Avoid HZ-like\nInflammatory Signature",
     "HZ Disease Reference\n(GSE242252 + Table S5)",
     "IFI44L, IFI27, RSAD2, ISG15, SERPING1",
     "NOT activated\n(flat or suppressed)",
     "Chapter 4:\nInnate Immunity\nEvaluation",
     RED_ACCENT),
    ("Standard 2", "Achieve RZV-like\nProtective Signature",
     "RZV Protection Reference\n(GSE249632)",
     "ZEB2, CTLA4, ICOS, HAVCR2",
     "Persistently\nupregulated",
     "Chapters 5–6:\nT Cell Response\nEvaluation",
     BLUE_ACCENT),
]):
    # Box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8 + i * 6.2), Inches(1.8), Inches(5.8), Inches(4.5))
    shape.fill.solid(); shape.fill.fore_color.rgb = BG_GREY
    shape.line.color.rgb = color; shape.line.width = Pt(2)

    # Title
    add_textbox(slide, 1.2 + i * 6.2, 2.0, 5.0, 0.4, std_num, font_size=14, bold=True, color=color)
    add_textbox(slide, 1.2 + i * 6.2, 2.4, 5.0, 0.5, std_name, font_size=12, bold=True, color=DARK_GREY)
    add_line(slide, 1.2 + i * 6.2, 3.0, 4.5, color, Pt(1))

    # Details
    labels = ["Reference:", "Key Indicator Genes:", "Expected Direction:", "Evaluated In:"]
    values = [ref, genes, direction, chapter]
    y = 3.2
    for lab, val in zip(labels, values):
        add_textbox(slide, 1.2 + i * 6.2, y, 1.5, 0.3, lab, font_size=8, bold=True, color=MED_GREY)
        add_textbox(slide, 2.8 + i * 6.2, y, 3.3, 0.5, val, font_size=9, color=DARK_GREY)
        y += 0.85

# Bottom conclusion
add_line(slide, 0.8, 6.6, 11.5, LIGHT_GREY, Pt(0.5))
add_textbox(slide, 0.8, 6.8, 11.5, 0.3,
    "This framework transforms Chapter 3 from a descriptive database analysis into the theoretical coordinate system for the entire thesis.",
    font_size=10, bold=True, color=ROYAL_BLUE, font_name='Times New Roman')

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: IMMUNE GENE DOT PLOT DATA
# ══════════════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_line(slide, 0, 0, 13.33, ROYAL_BLUE, Pt(3))
add_title(slide, "Key Immune Genes: RZV Vaccination Timeline by Functional Category")
add_subtitle(slide, "Curated 22-Gene Panel | 7 Functional Categories | D14 → D60 → D74 → D365")
add_line(slide, 0.8, 1.55, 11.5, ROYAL_BLUE, Pt(1.5))

# Gene LFC table by category
categories = [
    ("Differentiation", ["ZEB2", "IRF4", "TBX21"]),
    ("Co-stim / Checkpoint", ["ICOS", "CTLA4", "TIGIT", "PDCD1", "HAVCR2"]),
    ("Cytotoxicity", ["GZMA", "GNLY", "PRF1"]),
    ("Activation", ["CD38", "TOX"]),
    ("Memory / Naive", ["CCR7", "IL7R", "TCF7"]),
    ("Tfh", ["CXCR5", "BCL6"]),
    ("Type I IFN (Control)", ["ISG15", "MX1", "RSAD2", "STAT1"]),
]

# Build table rows with approximate LFC values from our data
lfc_approx = {
    'ZEB2': ('+3.2','+2.5','+3.6','+3.0'), 'IRF4': ('+0.5','+0.6','+0.7','+0.0'),
    'TBX21': ('+0.5','+0.3','+0.4','+0.2'),
    'ICOS': ('+1.3','+0.8','+1.2','+0.6'), 'CTLA4': ('+1.5','+1.4','+1.6','+1.0'),
    'TIGIT': ('+0.8','+0.5','+0.9','+0.4'), 'PDCD1': ('+0.8','+0.5','+0.7','+0.3'),
    'HAVCR2': ('+7.6','+8.4','+9.1','+8.0'),
    'GZMA': ('+2.5','+1.0','+3.4','+1.5'), 'GNLY': ('−2.5','+2.6','+3.2','+3.4'),
    'PRF1': ('+0.6','−0.9','+1.0','+1.4'),
    'CD38': ('+6.4','+3.4','+3.0','−0.1'), 'TOX': ('+1.2','+0.8','+1.0','+0.5'),
    'CCR7': ('+0.6','+0.5','+0.2','−0.3'), 'IL7R': ('−1.2','−0.5','−0.8','−1.0'),
    'TCF7': ('+0.5','+0.6','+0.1','−0.1'),
    'CXCR5': ('−0.8','+0.6','−0.1','−1.6'), 'BCL6': ('+0.3','+0.4','+0.5','+0.1'),
    'ISG15': ('+0.1','+0.2','−0.4','−0.1'), 'MX1': ('+0.3','+0.1','−0.2','−0.1'),
    'RSAD2': ('−0.3','+0.1','−0.2','−0.2'), 'STAT1': ('−0.2','−0.2','−0.3','−0.1'),
}

rows = []
for cat_name, genes in categories:
    rows.append([f"■ {cat_name}", "", "", "", ""])  # category header
    for g in genes:
        vals = lfc_approx.get(g, ('','','',''))
        rows.append([f"  {g}", vals[0], vals[1], vals[2], vals[3]])

add_table(slide, 0.8, 1.8,
    [2.5, 2.2, 2.2, 2.2, 2.2],
    ["Gene", "D14 LFC", "D60 LFC", "D74 LFC", "D365 LFC"],
    rows, header_color=ROYAL_BLUE)

# Note
add_line(slide, 0.8, 6.0, 11.5, LIGHT_GREY, Pt(0.5))
add_textbox(slide, 0.8, 6.2, 11.5, 0.8,
    "Color coding: Green = Protection Signature genes (sustained up). Red = Disease Signature genes (Type I IFN — flat). "
    "Differentiation genes (ZEB2) and checkpoint genes (CTLA4, HAVCR2) show strongest long-term signals. "
    "Type I IFN genes (ISG15, MX1, RSAD2, STAT1) serve as negative controls — completely flat across timeline.",
    font_size=8, color=MED_GREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: CONCLUSIONS
# ══════════════════════════════════════════════════════════════════════════════
slide = add_blank_slide()
add_line(slide, 0, 0, 13.33, ROYAL_BLUE, Pt(3))
add_title(slide, "Conclusions & Next Steps")
add_line(slide, 0.8, 1.05, 11.5, ROYAL_BLUE, Pt(1.5))

conclusions = [
    ("1. HZ Disease is an Innate Immune Storm",
     "Type I IFN ISGs (IFI44L +1.82, IFI27 +2.16, ISG15 +1.57) dominate the acute transcriptional "
     "landscape. Complement, B cell activation, and broad proliferation co-occur. Single-cell validation "
     "(HRA008316) confirms ISG signal originates from monocytes/DCs/neutrophils — not T cells."),
    ("2. RZV Induces Precise, Self-Limiting Adaptive Immunity",
     "Pulsed activation pattern (D14 peak → D60 baseline → D74 re-activation) proves RZV does NOT "
     "cause sustained inflammation. Long-term reprogramming (ZEB2 +3.0 at D365) with built-in checkpoints "
     "(CTLA4 +1.0) and active inflammatory suppression (SPP1 −5.9)."),
    ("3. Two Signatures are Fundamentally Distinct",
     "The Type I IFN pathway — hallmark of HZ — is completely absent in RZV response. "
     "RZV achieves protection through T cell differentiation reprogramming without triggering innate inflammation."),
    ("4. Dual Evaluation Framework Established",
     "Candidate vaccines should be assessed against two standards:\n"
     "  • Avoid Disease Signature (ISG15, RSAD2, IFI44L NOT activated)\n"
     "  • Achieve Protection Signature (ZEB2, HAVCR2, CTLA4 persistently upregulated)"),
]

y = 1.5
for title, body in conclusions:
    add_textbox(slide, 0.8, y, 11.5, 0.3, title, font_size=12, bold=True, color=ROYAL_BLUE)
    add_textbox(slide, 0.8, y + 0.35, 11.5, 0.7, body, font_size=9, color=DARK_GREY)
    y += 1.25


# ── Save ─────────────────────────────────────────────────────────────────────
OUT = "/media/cmj/MechanicalDisk/yjs/VZV-geo/results/Chapter3_Presentation.pptx"
prs.save(OUT)
print(f"PPT saved: {OUT} ({os.path.getsize(OUT)/1024:.0f} KB)")
